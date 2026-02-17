from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_figure_skating_presentation(output_path='figure_skating_presentation.pptx'):
    """Создает презентацию PowerPoint о топ-10 фигуристках"""

    # Данные для презентации (топ-10 фигуристок с их минусами)
    skaters = [
        {
            'name': 'Каори Сакамото (Япония)',
            'ranking': '1 место в рейтинге ISU',
            'weaknesses': [
                'Нестабильность в исполнении четверных прыжков',
                'Периодические ошибки на приземлениях',
                'Ограниченная артистическая выразительность в некоторых программах'
            ]
        },
        {
            'name': 'Алиса Лю (США)',
            'ranking': '2 место в рейтинге ISU',
            'weaknesses': [
                'Проблемы с консистентностью в течение сезона',
                'Сложности с удержанием концентрации в длительных программах',
                'Ограниченный набор четверных прыжков'
            ]
        },
        {
            'name': 'Хэ Ин (Китай)',
            'ranking': '3 место в рейтинге ISU',
            'weaknesses': [
                'Технические ошибки на сложных элементах',
                'Недостаточная скорость вращения',
                'Проблемы с выносливостью во второй половине программы'
            ]
        },
        {
            'name': 'Рика Кихира (Япония)',
            'ranking': '4 место в рейтинге ISU',
            'weaknesses': [
                'История травм, влияющая на стабильность',
                'Ограниченная вариативность в хореографии',
                'Сложности с адаптацией к новым правилам судейства'
            ]
        },
        {
            'name': 'Аделия Петросян (нейтральный статус)',
            'ranking': '5 место в рейтинге ISU',
            'weaknesses': [
                'Неопытность на крупных международных стартах',
                'Давление из-за нейтрального статуса',
                'Технические ошибки под стрессом'
            ]
        },
        {
            'name': 'Эмбер Гленн (США)',
            'ranking': '6 место в рейтинге ISU',
            'weaknesses': [
                'Нестабильность в короткой программе',
                'Ограниченная сложность вращений',
                'Проблемы с психологической устойчивостью'
            ]
        },
        {
            'name': 'Моне Тиба (Япония)',
            'ranking': '7 место в рейтинге ISU',
            'weaknesses': [
                'Недостаточная мощность прыжков',
                'Ограниченная артистическая составляющая',
                'Сложности с интерпретацией музыки'
            ]
        },
        {
            'name': 'Ким Чхэ Ён (Корея)',
            'ranking': '8 место в рейтинге ISU',
            'weaknesses': [
                'Технические ограничения в прыжках',
                'Нестабильность в исполнении каскадов',
                'Проблемы с сохранением энергии'
            ]
        },
        {
            'name': 'Николь Шотт (Германия)',
            'ranking': '9 место в рейтинге ISU',
            'weaknesses': [
                'Ограниченная техническая сложность',
                'Недостаточная скорость скольжения',
                'Консервативный подход к программам'
            ]
        },
        {
            'name': 'Мадлен Скизас (Канада)',
            'ranking': '10 место в рейтинге ISU',
            'weaknesses': [
                'Неопытность на топ-уровне',
                'Технические ошибки под давлением',
                'Ограниченная вариативность элементов'
            ]
        }
    ]

    # Создаем презентацию
    prs = Presentation()

    # Стиль для титульного слайда
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Топ-10 фигуристок\nОлимпиада 2026"
    subtitle.text = "Женское одиночное катание\nАнализ слабых сторон"

    # Слайд с введением
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "О презентации"
    tf = body_shape.text_frame
    tf.text = "Ключевая информация:"

    p = tf.add_paragraph()
    p.text = "• Олимпийские игры 2026 в Милане-Кортина д'Ампеццо"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "• Женское одиночное фигурное катание"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "• Топ-10 фигуристок по рейтингу ISU"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "• Фокус на слабых сторонах и минусах"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "• Актуальные данные сезона 2024-2025"
    p.level = 0

    # Создаем слайды для каждой фигуристки
    for i, skater in enumerate(skaters, 1):
        # Создаем новый слайд с заголовком и содержимым
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_shape.text = f"{i}. {skater['name']}"

        tf = body_shape.text_frame
        tf.text = f"Рейтинг: {skater['ranking']}"

        # Добавляем минусы
        p = tf.add_paragraph()
        p.text = "Основные минусы:"
        p.level = 0

        for weakness in skater['weaknesses']:
            p = tf.add_paragraph()
            p.text = f"• {weakness}"
            p.level = 1

    # Заключительный слайд
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = "Выводы"

    tf = body_shape.text_frame
    tf.text = "Ключевые наблюдения:"

    conclusions = [
        "1. Все топ-фигуристки имеют уязвимые места",
        "2. Техническая нестабильность - общая проблема",
        "3. Психологическая устойчивость критически важна",
        "4. Олимпийское давление усиливает слабые стороны",
        "5. Успех зависит от минимизации ошибок"
    ]

    for conclusion in conclusions:
        p = tf.add_paragraph()
        p.text = conclusion
        p.level = 0

    # Сохраняем презентацию
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    output = create_figure_skating_presentation()
    print(f"Презентация создана: {output}")